"""
Unified document reader — extracts clean text + metadata from any supported format.

Supported:  .pdf  .docx  .pptx  .xlsx  .xls  .csv
            .txt  .md  .py  .js  .ts  .html  .htm  .css
            .json  .xml  .yaml  .yml  .toml  .ini  .cfg  .env
            .java  .c  .cpp  .h  .cs  .go  .rs  .sh  .bat  .ps1  .sql

read_file(path)          -> FileContent
read_multiple_files(paths, workers=6) -> list[FileContent]
"""
from __future__ import annotations

import logging
import re
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

_MAX_BYTES = 512 * 1024   # 512 KB text cap per file
_MAX_ROWS  = 500          # Excel/CSV row preview cap
_MAX_COLS  = 50           # Excel/CSV column cap

_TEXT_EXTS = {
    ".txt", ".md", ".rst", ".py", ".js", ".ts", ".jsx", ".tsx",
    ".html", ".htm", ".css", ".json", ".xml", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".env", ".sh", ".bat", ".ps1",
    ".java", ".c", ".cpp", ".h", ".cs", ".go", ".rs", ".sql",
    ".log", ".gitignore", ".dockerfile", ".makefile",
}


# ─────────────────────────────────── data model ───────────────────────────────

@dataclass
class FileContent:
    path: str
    name: str
    format: str
    content: str = ""
    summary_meta: dict = field(default_factory=dict)
    error: str = ""

    @property
    def ok(self) -> bool:
        return not self.error

    def as_context(self, max_chars: int = 60_000) -> str:
        """Return a labelled block suitable for pasting into an AI prompt."""
        header = f"=== {self.name} ({self.format}) ==="
        if self.error:
            return f"{header}\n[ERROR: {self.error}]\n"
        meta = "  ".join(f"{k}: {v}" for k, v in self.summary_meta.items())
        body = self.content[:max_chars]
        tail = f"\n… (truncated, {len(self.content):,} chars total)" if len(self.content) > max_chars else ""
        return f"{header}\n[{meta}]\n\n{body}{tail}\n"


# ─────────────────────────────────── readers ──────────────────────────────────

def _read_text(p: Path) -> FileContent:
    try:
        raw = p.read_bytes()[:_MAX_BYTES]
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            text = raw.decode("latin-1", errors="replace")
        lines = text.splitlines()
        return FileContent(
            path=str(p), name=p.name, format=p.suffix.lstrip(".") or "text",
            content=text,
            summary_meta={"lines": len(lines), "chars": len(text)},
        )
    except OSError as exc:
        return FileContent(path=str(p), name=p.name, format="text", error=str(exc))


def _read_pdf(p: Path) -> FileContent:
    try:
        import pdfplumber
        pages_text: list[str] = []
        with pdfplumber.open(p) as pdf:
            n_pages = len(pdf.pages)
            for i, page in enumerate(pdf.pages, 1):
                t = page.extract_text() or ""
                if t.strip():
                    pages_text.append(f"[Page {i}]\n{t}")
        content = "\n\n".join(pages_text)
        return FileContent(
            path=str(p), name=p.name, format="pdf",
            content=content[:_MAX_BYTES],
            summary_meta={"pages": n_pages, "chars": len(content)},
        )
    except ImportError:
        return FileContent(path=str(p), name=p.name, format="pdf",
                           error="pdfplumber not installed: pip install pdfplumber")
    except Exception as exc:
        return FileContent(path=str(p), name=p.name, format="pdf", error=str(exc))


def _read_docx(p: Path) -> FileContent:
    try:
        from docx import Document
        doc = Document(str(p))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                prefix = f"[{para.style.name}] " if para.style and "Heading" in para.style.name else ""
                parts.append(prefix + para.text)
        # tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append("\t".join(c.text.strip() for c in row.cells))
            parts.append("\n".join(rows))
        content = "\n".join(parts)
        return FileContent(
            path=str(p), name=p.name, format="docx",
            content=content[:_MAX_BYTES],
            summary_meta={"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)},
        )
    except ImportError:
        return FileContent(path=str(p), name=p.name, format="docx",
                           error="python-docx not installed: pip install python-docx")
    except Exception as exc:
        return FileContent(path=str(p), name=p.name, format="docx", error=str(exc))


def _read_pptx(p: Path) -> FileContent:
    try:
        from pptx import Presentation
        prs = Presentation(str(p))
        slides_text: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
        content = "\n\n".join(slides_text)
        return FileContent(
            path=str(p), name=p.name, format="pptx",
            content=content[:_MAX_BYTES],
            summary_meta={"slides": len(prs.slides)},
        )
    except ImportError:
        return FileContent(path=str(p), name=p.name, format="pptx",
                           error="python-pptx not installed: pip install python-pptx")
    except Exception as exc:
        return FileContent(path=str(p), name=p.name, format="pptx", error=str(exc))


def _read_excel(p: Path) -> FileContent:
    try:
        import pandas as pd
        sheets: list[str] = []
        xl = pd.ExcelFile(str(p))
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name, nrows=_MAX_ROWS)
            df = df.iloc[:, :_MAX_COLS]
            sheets.append(
                f"[Sheet: {sheet_name}  {df.shape[0]} rows × {df.shape[1]} cols]\n"
                + df.to_string(index=False, max_rows=_MAX_ROWS, max_cols=_MAX_COLS)
            )
        content = "\n\n".join(sheets)
        return FileContent(
            path=str(p), name=p.name, format=p.suffix.lstrip("."),
            content=content[:_MAX_BYTES],
            summary_meta={"sheets": len(xl.sheet_names)},
        )
    except ImportError:
        return FileContent(path=str(p), name=p.name, format="excel",
                           error="pandas/openpyxl not installed: pip install pandas openpyxl")
    except Exception as exc:
        return FileContent(path=str(p), name=p.name, format="excel", error=str(exc))


def _read_csv(p: Path) -> FileContent:
    try:
        import pandas as pd
        df = pd.read_csv(str(p), nrows=_MAX_ROWS)
        df = df.iloc[:, :_MAX_COLS]
        content = df.to_string(index=False, max_rows=_MAX_ROWS, max_cols=_MAX_COLS)
        return FileContent(
            path=str(p), name=p.name, format="csv",
            content=content[:_MAX_BYTES],
            summary_meta={"rows": len(df), "cols": len(df.columns)},
        )
    except ImportError:
        return FileContent(path=str(p), name=p.name, format="csv",
                           error="pandas not installed: pip install pandas")
    except Exception as exc:
        return FileContent(path=str(p), name=p.name, format="csv", error=str(exc))


_READERS = {
    ".pdf":  _read_pdf,
    ".docx": _read_docx,
    ".docm": _read_docx,
    ".pptx": _read_pptx,
    ".pptm": _read_pptx,
    ".xlsx": _read_excel,
    ".xlsm": _read_excel,
    ".xls":  _read_excel,
    ".csv":  _read_csv,
}


# ─────────────────────────────────── public API ───────────────────────────────

def read_file(path: str | Path) -> FileContent:
    """Read a single file and return its content + metadata."""
    p = Path(path).expanduser().resolve()
    if not p.exists():
        return FileContent(path=str(p), name=p.name, format="", error="File not found")
    if not p.is_file():
        return FileContent(path=str(p), name=p.name, format="", error="Not a file")
    ext = p.suffix.lower()
    reader = _READERS.get(ext)
    if reader:
        return reader(p)
    if ext in _TEXT_EXTS or _is_text(p):
        return _read_text(p)
    return FileContent(
        path=str(p), name=p.name, format=ext.lstrip(".") or "binary",
        error=f"Unsupported format '{ext}'. Supported: pdf docx pptx xlsx csv + text/code files.",
    )


def read_multiple_files(paths: list[str | Path], workers: int = 6) -> list[FileContent]:
    """Read multiple files concurrently. Returns results in input order."""
    results: dict[int, FileContent] = {}
    with ThreadPoolExecutor(max_workers=min(workers, len(paths))) as pool:
        futures = {pool.submit(read_file, p): i for i, p in enumerate(paths)}
        for fut in as_completed(futures):
            idx = futures[fut]
            try:
                results[idx] = fut.result()
            except Exception as exc:
                p = paths[idx]
                results[idx] = FileContent(
                    path=str(p), name=Path(p).name, format="", error=str(exc)
                )
    return [results[i] for i in range(len(paths))]


def build_multi_file_context(paths: list[str | Path], max_total_chars: int = 120_000) -> str:
    """
    Read all paths and return a single context string ready to paste into a prompt.
    Trims per-file to stay within max_total_chars.
    """
    results = read_multiple_files(paths)
    per_file = max(8_000, max_total_chars // max(len(results), 1))
    blocks = [r.as_context(per_file) for r in results]
    header = f"[Context: {len(results)} file(s) attached]\n\n"
    return header + "\n".join(blocks)


def supported_formats() -> list[str]:
    exts = sorted(_READERS.keys()) + sorted(_TEXT_EXTS)
    return sorted(set(e.lstrip(".") for e in exts))


def _is_text(p: Path) -> bool:
    try:
        return b"\x00" not in p.read_bytes()[:1024]
    except OSError:
        return False
